#!/usr/bin/env python3
"""
replace_pptx_courses.py

Simple utility to replace placeholders in a PPTX using python-pptx.

Usage examples:
  # JSON list where each item maps to a slide by `slide_index` (0-based):
  python tools/replace_pptx_courses.py --pptx input.pptx --data data/courses.json --out output.pptx

  # CSV with a column `slide_index` and columns TITLE,DURATION,DESCRIPTION,image
  python tools/replace_pptx_courses.py --pptx input.pptx --data data/courses.csv --out output.pptx

Placeholder rules implemented by the script:
  - For slide-mapped items (have `slide_index`), the script will open that slide and
    replace text tokens in shapes: {{TITLE}}, {{DURATION}}, {{DESCRIPTION}}, etc.
  - If an item has `image` pointing to a local image path, any shape whose text
    contains the token {{IMAGE}} will be replaced by that image (keeps position/size).

Notes:
  - This is intentionally conservative. It won't try to intelligently match complex
    layouts. If your PPTX uses custom placeholder names, tell me and I can adapt the
    script to match them.
  - The script creates a backup of the input file if --backup is passed.
"""

import argparse
import json
import csv
import os
from pptx import Presentation
from pptx.util import Inches


def load_data(path):
    path = os.path.abspath(path)
    if path.lower().endswith('.json'):
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    elif path.lower().endswith('.csv'):
        with open(path, 'r', encoding='utf-8-sig') as f:
            reader = csv.DictReader(f)
            return list(reader)
    else:
        raise ValueError('Unsupported data format: ' + path)


def replace_tokens_in_shape(shape, mapping):
    """Replace tokens like {{TITLE}} in a shape's text_frame if present."""
    if not hasattr(shape, 'text'):
        return False
    text = shape.text
    original = text
    for k, v in mapping.items():
        token = '{{' + k.upper() + '}}'
        if token in text and v is not None:
            text = text.replace(token, str(v))
    if text != original:
        try:
            # Clear paragraphs then set simple text to preserve basic formatting
            shape.text = text
        except Exception:
            # Some shapes may not allow direct assignment; ignore
            pass
        return True
    return False


def replace_image_in_shape(slide, shape, image_path):
    """Replace a text-containing shape that has {{IMAGE}} token with an image.
       Keeps the original shape's left/top/width/height for placement.
    """
    if not os.path.exists(image_path):
        return False
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    try:
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        # remove placeholder shape
        sp = shape._element
        sp.getparent().remove(sp)
        return True
    except Exception:
        return False


def apply_slide_mapping(prs, item, report):
    # slide_index expected
    idx = item.get('slide_index')
    if idx is None:
        report.append(f"missing slide_index for item: {item}")
        return
    try:
        idx = int(idx)
    except Exception:
        report.append(f"invalid slide_index: {item.get('slide_index')}")
        return
    if idx < 0 or idx >= len(prs.slides):
        report.append(f"slide_index {idx} out of range (0..{len(prs.slides)-1})")
        return
    slide = prs.slides[idx]
    # create mapping of tokens from item keys
    mapping = {}
    for k, v in item.items():
        if k == 'slide_index' or v is None:
            continue
        mapping[k.upper()] = v

    replaced = 0
    # iterate shapes
    for shape in list(slide.shapes):
        try:
            # image handling: if shape has text containing {{IMAGE}} and item.image present
            if hasattr(shape, 'text') and '{{IMAGE}}' in (shape.text or '') and item.get('image'):
                ok = replace_image_in_shape(slide, shape, item.get('image'))
                if ok:
                    replaced += 1
                    continue
        except Exception:
            pass

        # text replacement
        if replace_tokens_in_shape(shape, mapping):
            replaced += 1

    report.append(f"slide {idx}: replaced {replaced} shapes")


def main():
    parser = argparse.ArgumentParser(description='Replace PPTX placeholders using JSON/CSV data')
    parser.add_argument('--pptx', required=True, help='input PPTX file')
    parser.add_argument('--data', required=True, help='JSON or CSV file with replacement data')
    parser.add_argument('--out', required=False, help='output PPTX file (default: output.pptx)', default='output.pptx')
    parser.add_argument('--backup', action='store_true', help='create a .bak of the input PPTX')
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print('Input PPTX not found:', args.pptx)
        return

    data = load_data(args.data)
    # normalize data to list
    if isinstance(data, dict):
        # if top-level dict with keys -> treat as single mapping to replace tokens globally
        data_items = [data]
    else:
        data_items = list(data)

    if args.backup:
        bak = args.pptx + '.bak'
        if not os.path.exists(bak):
            import shutil
            shutil.copy2(args.pptx, bak)
            print('Backup created:', bak)

    prs = Presentation(args.pptx)
    report = []

    # If items have slide_index, use slide-mapped mode
    has_slide_index = any(isinstance(it, dict) and 'slide_index' in it for it in data_items)

    if has_slide_index:
        for it in data_items:
            apply_slide_mapping(prs, it, report)
    else:
        # Global token mapping: replace tokens across all slides
        mapping = {}
        # use first item if list
        first = data_items[0] if data_items else {}
        for k, v in first.items():
            mapping[k.upper()] = v

        replaced = 0
        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if replace_tokens_in_shape(shape, mapping):
                    replaced += 1
        report.append(f'global replace: replaced tokens in {replaced} shapes')

    out_path = args.out
    prs.save(out_path)
    print('Saved output:', out_path)
    print('\nReport:')
    for r in report:
        print(' -', r)


if __name__ == '__main__':
    main()
